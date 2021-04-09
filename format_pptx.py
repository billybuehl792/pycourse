#!python
# read_pptx.py - read content from powerpoints

import json
import re
import sys
from pptx import Presentation
from docx import Document
from docx.shared import Pt, Inches
from docx.oxml.shared import OxmlElement
from lxml import etree


class Slide:

    def __init__(self, course, slide_id):
        self.course = course
        self.slide_id = slide_id
        self.slide = self.course.pres.slides.get(self.slide_id, default=None)
        self.slide_num = self.course.pres.slides.index(self.slide) + 1

    @property
    def slide_type(self):
        # return slide's type
        if self.slide_id in self.course.title_slides:
            return 'title'
        elif self.slide_id in self.course.section_header_slides:
            return 'section_header'
        elif self.slide_id in self.course.menu_slides:
            return 'menu'
        else:
            return 'standard'
    
    @property
    def slide_text(self):
        # return text in slide
        text_runs = []
        slide = self.slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(self.format_string(run.text))

        return text_runs
    
    @property
    def slide_notes(self):
        # return slide notes section (narration)
        if self.slide.has_notes_slide:
            slide_notes = self.slide.notes_slide.notes_text_frame.text
            return self.format_string(slide_notes)

        return ''

    @staticmethod
    def format_string(string):
        string = string.replace('\n', ' ')
        string = string.replace('’', "'")
        string = string.replace('‘', "'")
        string = string.replace('“', '"')
        string = string.replace('”', '"')
        string = string.replace('–', '-')
        string = string.replace('—', '-')
        string = string.replace('‐', '')
        string = string.replace('…', '...')
        string = string.replace('˚', ' degrees ')
        if not string.replace(' ', ''):
            return ''
        return string
    

class Course:
    
    def __init__(self, pptx_file):
        # self.file_id = file_id
        self.pptx_file = pptx_file
        self.pres = Presentation(self.pptx_file)
    
    @property
    def slide_ids(self):
        return [slide.slide_id for slide in self.pres.slides]

    @property
    def title_slides(self):
        title_layout_index = 0
        layout = self.pres.slide_master.slide_layouts[title_layout_index]
        return [slide.slide_id for slide in layout.used_by_slides]

    @property
    def section_header_slides(self):
        section_header_layout_index = 2
        layout = self.pres.slide_master.slide_layouts[section_header_layout_index]
        return [slide.slide_id for slide in layout.used_by_slides]
    
    @property
    def menu_slides(self):
        menu_slide_layout_index = 5
        layout = self.pres.slide_master.slide_layouts[menu_slide_layout_index]
        return [slide.slide_id for slide in layout.used_by_slides]

    @property
    def standard_slides(self):
        normal_slides = []
        for n, layout in enumerate(self.pres.slide_layouts):
            # if title slide or slide header
            if n in [0, 2, 5]:
                continue
            for slide in layout.used_by_slides:
                normal_slides.append(slide.slide_id)
        return normal_slides

    @property
    def course_title(self):
        # title slide - first text of first slide of powerpoint
        slide = Slide(self, self.slide_ids[0])
        return(slide.slide_text[0])

    @property
    def course_id(self):
        # title slide - last text of first slide of powerpoint
        slide = Slide(self, self.slide_ids[0])
        return(slide.slide_text[-1])
    
    @property
    def file_id(self):
        try:
            code = re.search('-(.*?)-', self.course_id).group(1)
            num = self.course_id[self.course_id.rfind('-')+1:]
            return f'{code}-{num}'
        except Exception:
            return 'FILE_ID_ERROR'

    @property
    def course(self):

        pptx = {}   # return dictionary of pptx

        pptx['course_title'] = self.course_title
        pptx['course_id'] = self.course_id
        pptx['study_guide_pdf'] = f'{self.course_id}_508.pdf'
        pptx['study_guide_print'] = f'{self.course_id}_StudyGuide.pdf'
        pptx['sections'] = [{'section_title': 'Introduction', 'slides': []}]


        for slide_id in self.slide_ids:
            slide = Slide(self, slide_id)
            text = slide.slide_text
            if slide.slide_type == 'section_header':
                section = {'section_title': text[0], 'slides': []}
                pptx['sections'].append(section)
            pptx['sections'][-1]['slides'].append(slide_id)

        return pptx

    def write_json(self):
        # write course dict to json
        filename = f'{self.course_id}.json'
        with open(filename, 'w', encoding='utf-8') as f:
            json.dump(self.course, f, ensure_ascii=False, indent=4)
        return filename

    def write_docx(self, *args):

        def prevent_doc_breakup(doc):
            # prevents document tables from spanning 2 pages
            tags = doc.element.xpath('//w:tr')
            rows = len(tags)
            for row in range(0,rows):
                tag = tags[row]                     # Specify which <w:r> tag you want
                child = OxmlElement('w:cantSplit')  # Create arbitrary tag
                tag.append(child)                   # Append in the new tag
            
        doc = Document()
        style = doc.styles['Normal']
        style.font.size, style.font.name = Pt(11), 'Calibri'
        doc.add_paragraph(f'{self.course_id} - Narration Script')
        doc.add_paragraph(self.course_title)
        table = doc.add_table(rows=0, cols=2)
        table.style = 'Table Grid'
        for item in self.slide_ids:
            slide = Slide(self, item)
            if slide.slide_notes:
                note = slide.slide_notes
                for func in args:
                    note = func(note, slide.slide_text, slide.slide_num)
                row_cells = table.add_row().cells
                row_cells[0].text = f'{self.file_id}_{slide.slide_num}'
                row_cells[1].text = note

        for cell in table.columns[0].cells:
            cell.width = Inches(0.5)
        for cell in table.columns[1].cells:
            cell.width = Inches(7)
        for section in doc.sections:
            section.top_margin = Inches(0.5)
            section.bottom_margin = Inches(0.5)
            section.left_margin = Inches(0.5)
            section.right_margin = Inches(0.5)

        prevent_doc_breakup(doc)
        doc_file = f'{self.course_id}_narration.docx'
        try:
            doc.save(doc_file)
        except PermissionError:
            print(f'{doc_file} exists! or invalid permissions!')
            return False
        return doc_file

    def write_xml(self):
        # write course narration XML
        # course xml
        the_course = etree.Element('theCourse')
        the_course.append(etree.Comment('Generated with pycourse!'))
        etree.SubElement(the_course, 'myCourseTitle').text = self.course_id
        etree.SubElement(the_course, 'studyGuidePDF').text = f'{self.course_id}.pdf'
        etree.SubElement(the_course, 'studyGuidePrint').text = f'{self.course_id}_StudyGuide.pdf'
        etree.SubElement(the_course, 'courseMenu').text = 'YES'

        for n, section in enumerate(self.course.get('sections')):
            # theSections elem
            the_sections = etree.SubElement(the_course, 'theSections', {'title': section.get('section_title')})
            
            file_num = 0
            for s in section.get('slides'):
                # sectionNumber elem
                section_num = etree.SubElement(the_sections, 'sectionNumber')

                # theFileToLoad elem
                the_file_to_load = etree.SubElement(section_num, 'theFileToLoad')
                the_file_to_load.text = f'{self.file_id}_s{n}_{file_num+1}.html'
                file_num += 1

                # closed caption elem
                closed_caption = etree.SubElement(section_num, 'closedCaptionText')
                slide = Slide(self, s)
                closed_caption.text = slide.slide_notes
        
        xml_file = f'{self.course_id}_narration.xml'
        tree_string = etree.tostring(the_course, pretty_print=True).decode('utf-8')
        with open(xml_file, 'w') as f:
            f.write(tree_string)

        return xml_file

    def write_txt(self, *args):
        txt_file = f'{self.course_id}_narration.txt'
        with open(txt_file, 'w') as f:
            for item in self.slide_ids:
                slide = Slide(self, item)
                if slide.slide_notes:
                    note = slide.slide_notes
                    for func in args:
                        note = func(note, slide.slide_text, slide.slide_num)
                    f.write(f'{note}\n\n')

        return txt_file


# custom slide filters
def filter_kc(slide_notes, slide_text, slide_num):
    # remove knowlege check notes slides
    if ' '.join(slide_text).lower().replace(' ', '').startswith('knowledgecheck'):
        print(f'Slide {str(slide_num)} | Knowledge Check skipped: {slide_notes}')
        return ''
    return slide_notes

def filter_ai(slide_notes, _, slide_num):
    # omit slide notes after "Aditional Information" in notes
    if 'Additional Information' in slide_notes:
        split_notes = slide_notes.split('Additional Information')
        if split_notes[0].replace('\n', ''):
            print(f'Slide {slide_num} | Additional Information skipped: ' + split_notes[1].replace('\n', ''))
            return split_notes[0].replace('\n', '')
        return ''
    return slide_notes


if __name__ == '__main__':
    # try:
    #     pres_file = sys.argv[1]
    # except IndexError:
    #     print('Provide arguments: <pptx file>')
    #     sys.exit()

    # pres_file = r'test_files\SMA-HQ-WBT-108.pptx'
    pres_file = r'test_files\SMA-SS-WBT-0013_RIDM.pptx'
    # pres_file = r'test_files\SMA-AS-WBT-101 NAMIS Refresher 11-6-2020.pptx'
    # pres_file = r'test_files\SMA-OV-WBT-132 Orion Capsure Recovery Case Study 03-04-21.pptx'
    # pres_file = r'test_files\SMA-OV-WBT-131_03-23-21.pptx'

    course = Course(pres_file)
    for slide in course.slide_ids:
        s = Slide(course, slide)
        print(s.slide_text)