#!python
# read_pptx.py - read content from powerpoints

import json
import os
import re
import sys
from pptx import Presentation
from docx import Document
from docx.shared import Pt, Inches
from docx.oxml.shared import OxmlElement
from lxml import etree

MENU_INDEX = 5
SECTION_HEADER_INDEX = 2
TITLE_INDEX = 0


class Course:

    def __init__(self, pptx_file, course_id=None, file_id=None, course_title=None):
        self.pptx_file = pptx_file
        self._course_id = course_id
        self._file_id = file_id
        self._course_title = course_title
        self.pres = Presentation(self.pptx_file)
    
    @property
    def slide_ids(self):
        return [slide.slide_id for slide in self.pres.slides]

    @property
    def title_slides(self):
        layout = self.pres.slide_master.slide_layouts[TITLE_INDEX]
        return [slide.slide_id for slide in layout.used_by_slides]

    @property
    def section_header_slides(self):
        layout = self.pres.slide_master.slide_layouts[SECTION_HEADER_INDEX]
        return [slide.slide_id for slide in layout.used_by_slides]
    
    @property
    def menu_slides(self):
        layout = self.pres.slide_master.slide_layouts[MENU_INDEX]
        return [slide.slide_id for slide in layout.used_by_slides]

    @property
    def standard_slides(self):
        normal_slides = []
        for n, layout in enumerate(self.pres.slide_layouts):
            # if title slide or slide header
            if n in [MENU_INDEX, SECTION_HEADER_INDEX, TITLE_INDEX]:
                continue
            for slide in layout.used_by_slides:
                normal_slides.append(slide.slide_id)
        return normal_slides

    @property
    def has_menu(self):
        # return True if len(menu_slides) > 0
        return True if len(self.menu_slides) > 0 else False

    @property
    def course_title(self):
        # title slide - first text of first slide of powerpoint
        if not self._course_title:
            slide_text = self.slide_text(1)
            return slide_text[0]
        else:
            return self._course_title

    @property
    def course_id(self):
        # title slide - last text of first slide of powerpoint
        if not self._course_id:
            slide_text = self.slide_text(1)
            return slide_text[-1]
        else:
            return self._course_id
    
    @property
    def file_id(self):
        # course_id - filtered text
        if not self._file_id:
            try:
                code = re.search('-(.*?)-', self.course_id).group(1)
                num = self.course_id[self.course_id.rfind('-')+1:]
                return f'{code}{num}'
            except Exception:
                return 'FILE_ID_ERROR'
        else:
            return self._file_id

    @property
    def course(self):
        # return dictionary of pptx course
        pptx = {}   

        pptx['course_title'] = self.course_title
        pptx['course_id'] = self.course_id
        pptx['study_guide_pdf'] = f'{self.course_id}_508.pdf'
        pptx['study_guide_print'] = f'{self.course_id}_StudyGuide.pdf'
        pptx['sections'] = [{'section_title': 'Introduction', 'slides': []}]

        for n, slide in enumerate(self.pres.slides):
            slide_text = self.slide_text(n+1)
            if self.slide_type(n+1) == 'section_header':
                section = {'section_title': slide_text[0], 'slides': []}
                pptx['sections'].append(section)
            pptx['sections'][-1]['slides'].append(slide.slide_id)

        return pptx
    
    def slide_text(self, slide_num=1):
        # return text present in slide
        text_runs = []
        slide = self.pres.slides[slide_num-1]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(self.format_string(run.text))

        return text_runs

    def slide_notes(self, slide_num=1):
        # return slide notes section (narration)
        slide = self.pres.slides[slide_num-1]
        if slide.has_notes_slide:
            slide_notes = slide.notes_slide.notes_text_frame.text
            return self.format_string(slide_notes)

        return ''

    def slide_type(self, slide_num=1):
        # return type of slide
        slide = self.pres.slides[slide_num-1]

        if slide.slide_id in self.title_slides:
            return 'title'
        elif slide.slide_id in self.section_header_slides:
            return 'section_header'
        elif slide.slide_id in self.menu_slides:
            return 'menu'
        else:
            return 'standard'

    def write_json(self):
        # write course dict to json
        filename = f'{self.course_id}.json'
        with open(filename, 'w', encoding='utf-8') as f:
            json.dump(self.course, f, ensure_ascii=False, indent=4)
            print(f'JSON written: {filename}')
        return filename

    def write_docx(self, *args):
        # write docx narration file

        def prevent_doc_breakup(doc):
            # prevents document tables from spanning 2 pages
            tags = doc.element.xpath('//w:tr')
            rows = len(tags)
            for row in range(0,rows):
                tag = tags[row]                     # Specify which <w:r> tag you want
                child = OxmlElement('w:cantSplit')  # Create arbitrary tag
                tag.append(child)                   # Append in the new tag
        
        def format_table(doc, table):
            for cell in table.columns[0].cells:
                cell.width = Inches(0.5)
            for cell in table.columns[1].cells:
                cell.width = Inches(7)
            for section in doc.sections:
                section.top_margin = Inches(0.5)
                section.bottom_margin = Inches(0.5)
                section.left_margin = Inches(0.5)
                section.right_margin = Inches(0.5)

        doc_file = f'{self.course_id}_narration.docx'
        doc = Document()
        style = doc.styles['Normal']
        style.font.size, style.font.name = Pt(11), 'Calibri'
        doc.add_paragraph(f'{self.course_id} - Narration Script')
        doc.add_paragraph(self.course_title)
        table = doc.add_table(rows=0, cols=2)
        table.style = 'Table Grid'

        for n, _ in enumerate(self.pres.slides):

            # slide information
            slide_num = n + 1
            slide_notes = self.slide_notes(slide_num)
            slide_text = self.slide_text(slide_num)
            slide_type = self.slide_type(slide_num)

            # skip main menu
            if slide.slide_type == 'menu' or slide.slide_type == 'section_header':
                continue
            
            # apply filters
            for func in args:
                slide_notes = func(slide_notes, slide_text, slide_num, slide_type)
            
            # add narration to table
            if slide_notes:
                row_cells = table.add_row().cells
                row_cells[0].text = f'{self.file_id}_{slide_num}'
                row_cells[1].text = slide_notes
            else:
                continue
        

        # format table cells
        format_table(doc, table)

        # prevent table cell from spanning 2 pages
        prevent_doc_breakup(doc)

        # save file to cwd
        try:
            doc.save(doc_file)
            print(f'DOCX written: {doc_file}')
        except Exception:
            print(f'{doc_file} is open! or invalid permissions!')
            return False

        return doc_file

    def write_xml(self, *args):
        # write XML narration file

        the_course = etree.Element('theCourse')
        the_course.append(etree.Comment('Generated with pycourse!'))
        etree.SubElement(the_course, 'myCourseTitle').text = self.course_title
        etree.SubElement(the_course, 'studyGuidePDF').text = f'{self.course_id}.pdf'
        etree.SubElement(the_course, 'studyGuidePrint').text = f'{self.course_id}_StudyGuide.pdf'

        # course menu
        menu = 'NO'
        if self.has_menu:
            menu = 'YES'

        etree.SubElement(the_course, 'courseMenu').text = menu

        for n, section in enumerate(self.course.get('sections')):
            
            # theSections elem
            the_sections = etree.SubElement(the_course, 'theSections', {'title': section.get('section_title')})
            
            file_num = 0
            for slide_id in section.get('slides'):
                
                # slide information
                slide_num = self.slide_ids.index(slide_id) + 1
                slide_notes = self.slide_notes(slide_num)
                slide_text = self.slide_text(slide_num)
                slide_type = self.slide_type(slide_num)

                # skip menus and section_headers
                if slide_type == 'menu' or slide_type == 'section_header':
                    continue

                # apply filters
                for func in args:
                    slide_notes = func(slide_notes, slide_text, slide_num, slide_type)

                # sectionNumber elem
                section_num = etree.SubElement(the_sections, 'sectionNumber')

                # theFileToLoad elem
                the_file_to_load = etree.SubElement(section_num, 'theFileToLoad')
                the_file_to_load.text = f'{self.file_id}_s{n}_{file_num+1}.html'

                # closed caption elem
                closed_caption = etree.SubElement(section_num, 'closedCaptionText')
                closed_caption.text = slide_notes

                # increment filenum
                file_num += 1

        xml_file = f'{self.course_id}_narration.xml'
        tree_string = etree.tostring(the_course, pretty_print=True).decode('utf-8')
        with open(xml_file, 'w') as f:
            f.write(tree_string)
            print(f'XML written: {xml_file}')

        return xml_file

    def write_txt(self, *args):
        # write slide info to TXT file
        txt_file = f'{self.course_id}_narration.txt'

        with open(txt_file, 'w') as f:
            f.write(f'COURSE TITLE: {self.course_title}\n')
            f.write(f'COURSE ID: {self.course_id}\n\n')

            for n, _ in enumerate(self.pres.slides):
                # slide information
                slide_num = n + 1
                slide_notes = self.slide_notes(slide_num)
                slide_text = self.slide_text(slide_num)
                slide_type = self.slide_type(slide_num)

                # apply filters
                for func in args:
                    slide_notes = func(slide_notes, slide_text, slide_num, slide_type)
                
                # write to file
                f.write(f'SLIDE:       {slide_num}\n')
                f.write(f'SLIDE TYPE:  {slide_type}\n')
                f.write(f'SLIDE TEXT:  {" | ".join(slide_text)}\n')
                f.write(f'SLIDE NOTES: {slide_notes}\n\n')

            print(f'TXT written: {txt_file}')

        return txt_file

    def __repr__(self):
        return f'Course({self.pptx_file}, {self.course_id}, {self.course_title})'

    @staticmethod
    def format_string(string):
        string = string.replace('\n', ' ')
        string = string.replace('’', "'")
        string = string.replace('‘', "'")
        string = string.replace('“', '"')
        string = string.replace('”', '"')
        string = string.replace('–', '-')
        string = string.replace('—', '-')
        string = string.replace('‐', '')
        string = string.replace('…', '...')
        string = string.replace('˚', ' degrees ')
        if not string.replace(' ', ''):
            return ''
        return string


# custom slide narration filters
def skip_kc(slide_notes, slide_text, slide_num, slide_type):
    # return '' if knowledge check narration
    del slide_type
    if ' '.join(slide_text).lower().replace(' ', '').startswith('knowledgecheck'):
        print(f'Slide {str(slide_num)} | Knowledge Check skipped: {slide_notes}')
        return ''
    return slide_notes

def filter_ai(slide_notes, slide_text, slide_num, slide_type):
    # omit slide notes after "Aditional Information" in notes
    del slide_text, slide_type
    if 'Additional Information' in slide_notes:
        split_notes = slide_notes.split('Additional Information')
        if split_notes[0].replace('\n', ''):
            print(f'Slide {slide_num} | Additional Information skipped: ' + split_notes[1].replace('\n', ''))
            return split_notes[0].replace('\n', '')
        return ''
    return slide_notes

def to_caps(slide_notes, slide_text, slide_num, slide_type):
    del slide_text, slide_num, slide_type
    return slide_notes.upper()

def skip_menu(slide_notes, slide_text, slide_num, slide_type):
    del slide_text, slide_num
    if slide_type == 'menu':
        return ''
    return slide_notes


if __name__ == '__main__':
    
    # get sys arguments
    l = len(sys.argv)
    if l <= 1:
        print('sys arguments: <pptx file> <course_id:optional> <course_title:optional> <file_id:optional>')
        pptx_file = input('pptx course file: ')
        course_id = input('course id ("enter" to skip): ') or None
        course_title = input('course title ("enter" to skip): ') or None
        file_id = input('file id ("enter" to skip): ') or None
    else:
        pptx_file = sys.argv[1]
        course_id = None
        file_id = None
        course_title = None
        if l >= 3:
            course_id = sys.argv[2]
            if l >= 4:
                course_title = sys.argv[3]
                if l >= 5:
                    file_id = sys.argv[4]
    
    # validate pptx file
    if not os.path.isfile(pptx_file) and not pptx_file.endswith('.pptx'):
        sys.exit(f'"{pptx_file}" is not a course file!')
    
    # course object
    course = Course(pptx_file, course_id=course_id, course_title=course_title, file_id=file_id)
    print(course)
    # pres_file_example = r'testFiles\SMA-HQ-WBT-108.pptx'
    # course = Course(pres_file_example, course_id='SMA-HQ-WBT-108')
